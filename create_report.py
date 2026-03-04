#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定义颜色
TITLE_COLOR = RGBColor(0, 0, 0)
ACCENT_COLOR = RGBColor(200, 16, 46)  # BBC红色
TEXT_COLOR = RGBColor(51, 51, 51)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    
    # 添加来源
    source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    source_frame = source_box.text_frame
    p = source_frame.paragraphs[0]
    p.text = "Data Source: BBC News | Report Date: March 2026"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(128, 128, 128)

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 添加分割线
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

# 第1页：标题页
add_title_slide(prs, "Iran Situation: Global Political and Economic Impact", "In-depth Analysis Report Based on BBC News")

# 第2页：执行摘要
add_content_slide(prs, "Executive Summary", [
    "• Military conflict between Iran, US, and Israel escalates globally",
    "• UK allows US to use military bases for defensive strikes on Iranian missile sites",
    "• Oil prices surge; global energy market faces pressure and uncertainty",
    "• Geopolitical tensions in Middle East intensify; impact on trade and investment",
    "• Western nations condemn Iran's indiscriminate and disproportionate strikes"
])

# 第3页：政治影响
add_content_slide(prs, "Political Impact Analysis", [
    "• International alliance restructuring: US, Israel strengthen military cooperation",
    "• UK policy shift: Permits US use of RAF Fairford and Diego Garcia bases",
    "• Iran regional influence: Missile attacks on Gulf states, Iraq, Bahrain",
    "• Western stance: Emphasizes self-defense rights and international law legitimacy",
    "• Nuclear concerns: Heightened worries about Iran's nuclear weapons program"
])

# 第4页：经济影响
add_content_slide(prs, "Economic Impact Analysis", [
    "• Oil price volatility: Ships attacked near Strait of Hormuz; prices surge",
    "• Energy security: Global energy prices face sustained upward pressure",
    "• Trade disruption: Middle East air transport halted; business activities affected",
    "• Investment risk: Geopolitical risks reduce investor confidence",
    "• UK interests: 200,000 British citizens in Middle East face risks"
])

# 第5页：风险展望
add_content_slide(prs, "Risk Outlook and Prospects", [
    "• Escalation risk: Conflict may expand into broader regional warfare",
    "• Economic recession: Rising energy prices may trigger global slowdown",
    "• Humanitarian crisis: Large-scale conflict could cause refugee crisis",
    "• Nuclear threat: Uncertainty in Iran's nuclear program increases global security risks",
    "• International division: Western nations diverge from other powers on Iran issue"
])

# 第6页：表格 - 影响对比
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Impact Comparison Table"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR

# 添加表格
rows, cols = 4, 3
left = Inches(0.7)
top = Inches(1.2)
width = Inches(8.6)
height = Inches(5.5)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 设置列宽
table_shape.columns[0].width = Inches(2.5)
table_shape.columns[1].width = Inches(3)
table_shape.columns[2].width = Inches(3.1)

# 填充表格
data = [
    ["Dimension", "Political Impact", "Economic Impact"],
    ["Energy", "Strategic control of oil routes", "Oil prices up 5-10%"],
    ["Trade", "Regional instability", "Shipping delays, higher costs"],
    ["Investment", "Reduced confidence", "Capital flight from region"]
]

for i, row_data in enumerate(data):
    for j, cell_data in enumerate(row_data):
        cell = table_shape.cell(i, j)
        cell.text = cell_data
        
        # 设置表头样式
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_COLOR
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(14)
        else:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

# 第7页：BBC信息来源
add_content_slide(prs, "Primary BBC News Sources", [
    "1. 'UK will allow US to use bases to strike Iranian missile sites, PM says'",
    "   https://www.bbc.com/news/articles/cqj9g11p1ezo",
    "",
    "2. 'Oil prices jump after ships attacked near Strait of Hormuz'",
    "   https://www.bbc.com/news/articles/c75evve6l63o",
    "",
    "3. 'UK mounts operation to support thousands of Britons in Middle East'",
    "   https://www.bbc.com/news/articles/cdjm8pn0ld8o"
])

# 保存文件
output_path = "/Users/federerx/Desktop/ui-tars/iran_bbc_report.pptx"
prs.save(output_path)
print(f"PPTX Report created successfully: {output_path}")
